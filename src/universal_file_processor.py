"""
Universal File Format Processor
Handles Excel, PowerPoint, images, and all common file formats with intelligent processing.
"""

import os
import json
import logging
import tempfile
import mimetypes
from typing import Dict, List, Any, Optional, Tuple, Union
from pathlib import Path
import pandas as pd
import numpy as np
from dataclasses import dataclass

# Excel processing
try:
    import openpyxl
    import xlrd
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False
    logging.warning("Excel libraries not available. Install: pip install openpyxl xlrd")

# PowerPoint processing
try:
    from pptx import Presentation
    from pptx.shapes.picture import Picture
    POWERPOINT_AVAILABLE = True
except ImportError:
    POWERPOINT_AVAILABLE = False
    logging.warning("PowerPoint library not available. Install: pip install python-pptx")

# Image processing
try:
    from PIL import Image, ImageEnhance
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False
    logging.warning("PIL not available. Install: pip install Pillow")

try:
    from wand.image import Image as WandImage
    WAND_AVAILABLE = True
except ImportError:
    WAND_AVAILABLE = False
    logging.warning("Wand not available. Install: pip install Wand")

# Text extraction from various formats
try:
    import magic
    MAGIC_AVAILABLE = True
except ImportError:
    MAGIC_AVAILABLE = False
    logging.warning("python-magic not available. Install: pip install python-magic")

try:
    import xml.etree.ElementTree as ET
    import yaml
    import csv
    import html2text
    STRUCTURED_FORMATS_AVAILABLE = True
except ImportError:
    STRUCTURED_FORMATS_AVAILABLE = False
    logging.warning("Some structured format libraries not available")

# Ebook processing
try:
    import ebooklib
    from ebooklib import epub
    EBOOK_AVAILABLE = True
except ImportError:
    EBOOK_AVAILABLE = False
    logging.warning("Ebook library not available. Install: pip install ebooklib")

from langchain.schema import Document
from .config import Config
from .multimodal_rag import MultiModalElement


@dataclass
class FileProcessingResult:
    """Result of file processing with metadata."""
    documents: List[Document]
    multimodal_elements: List[MultiModalElement]
    file_type: str
    processing_method: str
    success: bool
    error: Optional[str] = None
    metadata: Dict[str, Any] = None


class UniversalFileProcessor:
    """Universal processor for all supported file formats."""
    
    def __init__(self, config: Optional[Config] = None):
        self.config = config or Config()
        
        # Import multimodal processors
        from .multimodal_rag import TableProcessor, ImageProcessor
        self.table_processor = TableProcessor(config)
        self.image_processor = ImageProcessor(config)
    
    def process_file(self, file_path: str) -> FileProcessingResult:
        """Process any supported file format."""
        try:
            path = Path(file_path)
            if not path.exists():
                return FileProcessingResult(
                    documents=[], multimodal_elements=[], file_type="unknown",
                    processing_method="none", success=False, error="File not found"
                )
            
            extension = path.suffix.lower()
            
            # Route to appropriate processor based on file extension
            if extension in [".xlsx", ".xls"]:
                return self._process_excel(file_path)
            elif extension == ".csv":
                return self._process_csv(file_path)
            elif extension in [".pptx", ".ppt"]:
                return self._process_powerpoint(file_path)
            elif extension in [".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff", ".webp"]:
                return self._process_image(file_path)
            elif extension == ".svg":
                return self._process_svg(file_path)
            elif extension == ".json":
                return self._process_json(file_path)
            elif extension in [".xml"]:
                return self._process_xml(file_path)
            elif extension in [".yaml", ".yml"]:
                return self._process_yaml(file_path)
            elif extension in [".html", ".htm"]:
                return self._process_html(file_path)
            elif extension == ".rtf":
                return self._process_rtf(file_path)
            elif extension in [".epub", ".mobi"]:
                return self._process_ebook(file_path)
            else:
                # Try to detect file type using magic if available
                return self._process_generic(file_path)
                
        except Exception as e:
            logging.error(f"File processing failed for {file_path}: {e}")
            return FileProcessingResult(
                documents=[], multimodal_elements=[], file_type="unknown",
                processing_method="error", success=False, error=str(e)
            )
    
    def _process_excel(self, file_path: str) -> FileProcessingResult:
        """Process Excel files (.xlsx, .xls)."""
        if not EXCEL_AVAILABLE:
            return FileProcessingResult(
                documents=[], multimodal_elements=[], file_type="excel",
                processing_method="unavailable", success=False, 
                error="Excel processing libraries not available"
            )
        
        try:
            path = Path(file_path)
            documents = []
            multimodal_elements = []
            
            # Read all sheets from Excel file
            with pd.ExcelFile(file_path) as excel_file:
                sheet_names = excel_file.sheet_names
                
                for sheet_name in sheet_names:
                    try:
                        # Read sheet as DataFrame
                        df = pd.read_excel(file_path, sheet_name=sheet_name, header=0)
                        
                        if df.empty:
                            continue
                        
                        # Clean the DataFrame
                        df = self._clean_dataframe(df)
                        
                        if not df.empty:
                            # Create document with sheet text content
                            sheet_text = self._dataframe_to_text(df, sheet_name)
                            doc = Document(
                                page_content=sheet_text,
                                metadata={
                                    "source": file_path,
                                    "filename": path.name,
                                    "file_type": path.suffix.lower(),
                                    "sheet_name": sheet_name,
                                    "processing_method": "pandas_excel",
                                    "rows": len(df),
                                    "columns": len(df.columns)
                                }
                            )
                            documents.append(doc)
                            
                            # Create multimodal table element
                            table_element = MultiModalElement(
                                element_id=f"excel_sheet_{sheet_name}_{hash(file_path) % 10000}",
                                element_type="table",
                                content=df,
                                metadata={
                                    "source": file_path,
                                    "sheet_name": sheet_name,
                                    "extraction_method": "pandas_excel",
                                    "file_type": "excel"
                                },
                                text_description=f"Excel sheet '{sheet_name}' with {len(df)} rows and {len(df.columns)} columns",
                                structured_data=df.to_dict(),
                                confidence_score=0.95,
                                processing_method="excel_pandas"
                            )
                            multimodal_elements.append(table_element)
                    
                    except Exception as sheet_error:
                        logging.warning(f"Failed to process sheet '{sheet_name}': {sheet_error}")
                        continue
            
            return FileProcessingResult(
                documents=documents,
                multimodal_elements=multimodal_elements,
                file_type="excel",
                processing_method="pandas_excel",
                success=True,
                metadata={
                    "sheets_processed": len(documents),
                    "total_sheets": len(sheet_names)
                }
            )
            
        except Exception as e:
            return FileProcessingResult(
                documents=[], multimodal_elements=[], file_type="excel",
                processing_method="pandas_excel", success=False, error=str(e)
            )
    
    def _process_csv(self, file_path: str) -> FileProcessingResult:
        """Process CSV files."""
        try:
            path = Path(file_path)
            
            # Try to read CSV with different parameters
            encodings = ['utf-8', 'latin1', 'cp1252']
            separators = [',', ';', '\t', '|']
            
            df = None
            used_encoding = None
            used_separator = None
            
            for encoding in encodings:
                for sep in separators:
                    try:
                        df = pd.read_csv(file_path, encoding=encoding, separator=sep, header=0)
                        if len(df.columns) > 1 and len(df) > 0:  # Reasonable CSV
                            used_encoding = encoding
                            used_separator = sep
                            break
                    except:
                        continue
                if df is not None:
                    break
            
            if df is None:
                # Fallback: read as plain text
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
                return FileProcessingResult(
                    documents=[Document(
                        page_content=content,
                        metadata={
                            "source": file_path,
                            "filename": path.name,
                            "file_type": path.suffix.lower(),
                            "processing_method": "text_fallback"
                        }
                    )],
                    multimodal_elements=[],
                    file_type="csv",
                    processing_method="text_fallback",
                    success=True
                )
            
            # Clean the DataFrame
            df = self._clean_dataframe(df)
            
            # Create text content
            csv_text = self._dataframe_to_text(df, "CSV Data")
            doc = Document(
                page_content=csv_text,
                metadata={
                    "source": file_path,
                    "filename": path.name,
                    "file_type": path.suffix.lower(),
                    "processing_method": "pandas_csv",
                    "encoding": used_encoding,
                    "separator": used_separator,
                    "rows": len(df),
                    "columns": len(df.columns)
                }
            )
            
            # Create table element
            table_element = MultiModalElement(
                element_id=f"csv_table_{hash(file_path) % 10000}",
                element_type="table",
                content=df,
                metadata={
                    "source": file_path,
                    "extraction_method": "pandas_csv",
                    "file_type": "csv",
                    "encoding": used_encoding,
                    "separator": used_separator
                },
                text_description=f"CSV data with {len(df)} rows and {len(df.columns)} columns",
                structured_data=df.to_dict(),
                confidence_score=0.9,
                processing_method="csv_pandas"
            )
            
            return FileProcessingResult(
                documents=[doc],
                multimodal_elements=[table_element],
                file_type="csv",
                processing_method="pandas_csv",
                success=True
            )
            
        except Exception as e:
            return FileProcessingResult(
                documents=[], multimodal_elements=[], file_type="csv",
                processing_method="pandas_csv", success=False, error=str(e)
            )
    
    def _process_powerpoint(self, file_path: str) -> FileProcessingResult:
        """Process PowerPoint files (.pptx, .ppt)."""
        if not POWERPOINT_AVAILABLE:
            return FileProcessingResult(
                documents=[], multimodal_elements=[], file_type="powerpoint",
                processing_method="unavailable", success=False,
                error="PowerPoint processing library not available"
            )
        
        try:
            path = Path(file_path)
            documents = []
            multimodal_elements = []
            
            # Only .pptx is supported by python-pptx
            if path.suffix.lower() == ".ppt":
                return FileProcessingResult(
                    documents=[], multimodal_elements=[], file_type="powerpoint",
                    processing_method="unsupported_format", success=False,
                    error="Legacy .ppt format not supported, use .pptx"
                )
            
            prs = Presentation(file_path)
            
            for i, slide in enumerate(prs.slides):
                slide_text = []
                slide_images = []
                slide_tables = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                    
                    # Extract tables
                    if shape.has_table:
                        table_data = []
                        for row in shape.table.rows:
                            row_data = []
                            for cell in row.cells:
                                row_data.append(cell.text.strip())
                            table_data.append(row_data)
                        
                        if table_data and len(table_data) > 1:
                            # Convert to DataFrame
                            try:
                                df = pd.DataFrame(table_data[1:], columns=table_data[0])
                                df = self._clean_dataframe(df)
                                
                                if not df.empty:
                                    table_element = MultiModalElement(
                                        element_id=f"ppt_table_slide{i+1}_{len(slide_tables)}",
                                        element_type="table",
                                        content=df,
                                        metadata={
                                            "source": file_path,
                                            "slide_number": i + 1,
                                            "extraction_method": "python_pptx",
                                            "file_type": "powerpoint"
                                        },
                                        text_description=f"Table from slide {i+1} with {len(df)} rows and {len(df.columns)} columns",
                                        structured_data=df.to_dict(),
                                        confidence_score=0.85,
                                        processing_method="powerpoint_table"
                                    )
                                    multimodal_elements.append(table_element)
                                    slide_tables.append(f"Table: {df.shape[0]} rows × {df.shape[1]} columns")
                            except Exception as table_error:
                                logging.warning(f"Failed to process table on slide {i+1}: {table_error}")
                    
                    # Note: Image extraction from PowerPoint is complex and would require additional work
                    # For now, we'll note that images exist
                    if isinstance(shape, Picture):
                        slide_images.append("Image detected")
                
                # Create document for slide
                if slide_text or slide_tables or slide_images:
                    slide_content = f"Slide {i+1}:\n"
                    if slide_text:
                        slide_content += "\nText content:\n" + "\n".join(slide_text)
                    if slide_tables:
                        slide_content += "\nTables:\n" + "\n".join(slide_tables)
                    if slide_images:
                        slide_content += f"\nImages: {len(slide_images)} image(s) detected"
                    
                    doc = Document(
                        page_content=slide_content,
                        metadata={
                            "source": file_path,
                            "filename": path.name,
                            "file_type": path.suffix.lower(),
                            "slide_number": i + 1,
                            "processing_method": "python_pptx",
                            "has_images": len(slide_images) > 0,
                            "has_tables": len(slide_tables) > 0
                        }
                    )
                    documents.append(doc)
            
            return FileProcessingResult(
                documents=documents,
                multimodal_elements=multimodal_elements,
                file_type="powerpoint",
                processing_method="python_pptx",
                success=True,
                metadata={
                    "total_slides": len(prs.slides),
                    "processed_slides": len(documents)
                }
            )
            
        except Exception as e:
            return FileProcessingResult(
                documents=[], multimodal_elements=[], file_type="powerpoint",
                processing_method="python_pptx", success=False, error=str(e)
            )
    
    def _process_image(self, file_path: str) -> FileProcessingResult:
        """Process image files (JPG, PNG, GIF, etc.)."""
        try:
            path = Path(file_path)
            
            # Use existing image processor
            if self.image_processor:
                image_element = self.image_processor.process_image(file_path)
                
                if image_element:
                    # Create text document with image description
                    doc = Document(
                        page_content=f"Image: {path.name}\n\nDescription: {image_element.text_description}",
                        metadata={
                            "source": file_path,
                            "filename": path.name,
                            "file_type": path.suffix.lower(),
                            "processing_method": "multimodal_image_processor",
                            "image_analyzed": True
                        }
                    )
                    
                    return FileProcessingResult(
                        documents=[doc],
                        multimodal_elements=[image_element],
                        file_type="image",
                        processing_method="multimodal_image_processor",
                        success=True
                    )
            
            # Fallback: basic image info
            if PIL_AVAILABLE:
                with Image.open(file_path) as img:
                    width, height = img.size
                    format_name = img.format
                    mode = img.mode
                    
                    doc = Document(
                        page_content=f"Image: {path.name}\nDimensions: {width}x{height}\nFormat: {format_name}\nColor mode: {mode}",
                        metadata={
                            "source": file_path,
                            "filename": path.name,
                            "file_type": path.suffix.lower(),
                            "processing_method": "pil_basic",
                            "width": width,
                            "height": height,
                            "format": format_name,
                            "mode": mode
                        }
                    )
                    
                    return FileProcessingResult(
                        documents=[doc],
                        multimodal_elements=[],
                        file_type="image",
                        processing_method="pil_basic",
                        success=True
                    )
            
            # Final fallback
            return FileProcessingResult(
                documents=[Document(
                    page_content=f"Image file: {path.name}",
                    metadata={"source": file_path, "filename": path.name, "file_type": path.suffix.lower()}
                )],
                multimodal_elements=[],
                file_type="image",
                processing_method="basic_info",
                success=True
            )
            
        except Exception as e:
            return FileProcessingResult(
                documents=[], multimodal_elements=[], file_type="image",
                processing_method="image_processing", success=False, error=str(e)
            )
    
    def _process_json(self, file_path: str) -> FileProcessingResult:
        """Process JSON files."""
        try:
            path = Path(file_path)
            
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            # Convert JSON to readable text
            json_text = self._json_to_text(data)
            
            doc = Document(
                page_content=json_text,
                metadata={
                    "source": file_path,
                    "filename": path.name,
                    "file_type": path.suffix.lower(),
                    "processing_method": "json_parser",
                    "data_type": type(data).__name__
                }
            )
            
            # If JSON contains tabular data, create table elements
            multimodal_elements = []
            if isinstance(data, list) and len(data) > 0 and isinstance(data[0], dict):
                # Looks like a table
                try:
                    df = pd.DataFrame(data)
                    df = self._clean_dataframe(df)
                    
                    if not df.empty:
                        table_element = MultiModalElement(
                            element_id=f"json_table_{hash(file_path) % 10000}",
                            element_type="table",
                            content=df,
                            metadata={
                                "source": file_path,
                                "extraction_method": "json_to_dataframe",
                                "file_type": "json"
                            },
                            text_description=f"JSON data converted to table with {len(df)} rows and {len(df.columns)} columns",
                            structured_data=df.to_dict(),
                            confidence_score=0.9,
                            processing_method="json_dataframe"
                        )
                        multimodal_elements.append(table_element)
                except Exception as table_error:
                    logging.warning(f"Failed to convert JSON to table: {table_error}")
            
            return FileProcessingResult(
                documents=[doc],
                multimodal_elements=multimodal_elements,
                file_type="json",
                processing_method="json_parser",
                success=True
            )
            
        except Exception as e:
            return FileProcessingResult(
                documents=[], multimodal_elements=[], file_type="json",
                processing_method="json_parser", success=False, error=str(e)
            )
    
    def _process_xml(self, file_path: str) -> FileProcessingResult:
        """Process XML files."""
        try:
            path = Path(file_path)
            
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            # Parse XML
            try:
                root = ET.fromstring(content)
                xml_text = self._xml_to_text(root)
            except ET.ParseError:
                # If XML parsing fails, treat as text
                xml_text = content
            
            doc = Document(
                page_content=xml_text,
                metadata={
                    "source": file_path,
                    "filename": path.name,
                    "file_type": path.suffix.lower(),
                    "processing_method": "xml_parser"
                }
            )
            
            return FileProcessingResult(
                documents=[doc],
                multimodal_elements=[],
                file_type="xml",
                processing_method="xml_parser",
                success=True
            )
            
        except Exception as e:
            return FileProcessingResult(
                documents=[], multimodal_elements=[], file_type="xml",
                processing_method="xml_parser", success=False, error=str(e)
            )
    
    def _process_yaml(self, file_path: str) -> FileProcessingResult:
        """Process YAML files."""
        try:
            path = Path(file_path)
            
            with open(file_path, 'r', encoding='utf-8') as f:
                data = yaml.safe_load(f)
            
            # Convert YAML to readable text
            yaml_text = self._yaml_to_text(data)
            
            doc = Document(
                page_content=yaml_text,
                metadata={
                    "source": file_path,
                    "filename": path.name,
                    "file_type": path.suffix.lower(),
                    "processing_method": "yaml_parser"
                }
            )
            
            return FileProcessingResult(
                documents=[doc],
                multimodal_elements=[],
                file_type="yaml",
                processing_method="yaml_parser",
                success=True
            )
            
        except Exception as e:
            return FileProcessingResult(
                documents=[], multimodal_elements=[], file_type="yaml",
                processing_method="yaml_parser", success=False, error=str(e)
            )
    
    def _process_html(self, file_path: str) -> FileProcessingResult:
        """Process HTML files."""
        try:
            path = Path(file_path)
            
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                html_content = f.read()
            
            # Convert HTML to text
            try:
                import html2text
                h = html2text.HTML2Text()
                h.ignore_links = False
                h.ignore_images = False
                text_content = h.handle(html_content)
            except ImportError:
                # Fallback: basic HTML tag removal
                import re
                text_content = re.sub('<[^<]+?>', '', html_content)
            
            doc = Document(
                page_content=text_content,
                metadata={
                    "source": file_path,
                    "filename": path.name,
                    "file_type": path.suffix.lower(),
                    "processing_method": "html2text"
                }
            )
            
            # Extract tables from HTML using existing table processor
            multimodal_elements = []
            if self.table_processor:
                try:
                    table_elements = self.table_processor.extract_tables_from_html(html_content)
                    multimodal_elements.extend(table_elements)
                except Exception as table_error:
                    logging.warning(f"HTML table extraction failed: {table_error}")
            
            return FileProcessingResult(
                documents=[doc],
                multimodal_elements=multimodal_elements,
                file_type="html",
                processing_method="html2text",
                success=True
            )
            
        except Exception as e:
            return FileProcessingResult(
                documents=[], multimodal_elements=[], file_type="html",
                processing_method="html2text", success=False, error=str(e)
            )
    
    def _process_rtf(self, file_path: str) -> FileProcessingResult:
        """Process RTF files."""
        try:
            path = Path(file_path)
            
            # RTF processing is complex, for now read as text with encoding handling
            encodings = ['utf-8', 'latin1', 'cp1252']
            content = None
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        content = f.read()
                    break
                except UnicodeDecodeError:
                    continue
            
            if content is None:
                raise ValueError("Could not decode RTF file")
            
            # Basic RTF tag removal (very basic)
            import re
            text_content = re.sub(r'\\[a-z]+\d*', '', content)
            text_content = re.sub(r'[{}]', '', text_content)
            text_content = '\n'.join(line.strip() for line in text_content.split('\n') if line.strip())
            
            doc = Document(
                page_content=text_content,
                metadata={
                    "source": file_path,
                    "filename": path.name,
                    "file_type": path.suffix.lower(),
                    "processing_method": "basic_rtf"
                }
            )
            
            return FileProcessingResult(
                documents=[doc],
                multimodal_elements=[],
                file_type="rtf",
                processing_method="basic_rtf",
                success=True
            )
            
        except Exception as e:
            return FileProcessingResult(
                documents=[], multimodal_elements=[], file_type="rtf",
                processing_method="basic_rtf", success=False, error=str(e)
            )
    
    def _process_svg(self, file_path: str) -> FileProcessingResult:
        """Process SVG files."""
        try:
            path = Path(file_path)
            
            with open(file_path, 'r', encoding='utf-8') as f:
                svg_content = f.read()
            
            # Extract text from SVG
            import re
            text_elements = re.findall(r'<text[^>]*>(.*?)</text>', svg_content, re.IGNORECASE)
            extracted_text = '\n'.join(text_elements) if text_elements else "SVG graphic content"
            
            doc = Document(
                page_content=f"SVG file: {path.name}\n\nExtracted text:\n{extracted_text}",
                metadata={
                    "source": file_path,
                    "filename": path.name,
                    "file_type": path.suffix.lower(),
                    "processing_method": "svg_text_extraction",
                    "text_elements_found": len(text_elements)
                }
            )
            
            return FileProcessingResult(
                documents=[doc],
                multimodal_elements=[],
                file_type="svg",
                processing_method="svg_text_extraction",
                success=True
            )
            
        except Exception as e:
            return FileProcessingResult(
                documents=[], multimodal_elements=[], file_type="svg",
                processing_method="svg_text_extraction", success=False, error=str(e)
            )
    
    def _process_ebook(self, file_path: str) -> FileProcessingResult:
        """Process ebook files (EPUB, MOBI)."""
        if not EBOOK_AVAILABLE:
            return FileProcessingResult(
                documents=[], multimodal_elements=[], file_type="ebook",
                processing_method="unavailable", success=False,
                error="Ebook processing library not available"
            )
        
        try:
            path = Path(file_path)
            documents = []
            
            if path.suffix.lower() == ".epub":
                book = epub.read_epub(file_path)
                
                for item in book.get_items():
                    if item.get_type() == ebooklib.ITEM_DOCUMENT:
                        content = item.get_content().decode('utf-8')
                        
                        # Convert HTML to text
                        try:
                            import html2text
                            h = html2text.HTML2Text()
                            text_content = h.handle(content)
                        except ImportError:
                            import re
                            text_content = re.sub('<[^<]+?>', '', content)
                        
                        if text_content.strip():
                            doc = Document(
                                page_content=text_content,
                                metadata={
                                    "source": file_path,
                                    "filename": path.name,
                                    "file_type": path.suffix.lower(),
                                    "processing_method": "ebooklib_epub",
                                    "item_id": item.get_id(),
                                    "item_name": item.get_name()
                                }
                            )
                            documents.append(doc)
            
            else:
                # MOBI processing would require additional libraries
                return FileProcessingResult(
                    documents=[], multimodal_elements=[], file_type="ebook",
                    processing_method="unsupported_format", success=False,
                    error="MOBI format not yet supported"
                )
            
            return FileProcessingResult(
                documents=documents,
                multimodal_elements=[],
                file_type="ebook",
                processing_method="ebooklib_epub",
                success=True,
                metadata={"chapters_processed": len(documents)}
            )
            
        except Exception as e:
            return FileProcessingResult(
                documents=[], multimodal_elements=[], file_type="ebook",
                processing_method="ebooklib_epub", success=False, error=str(e)
            )
    
    def _process_generic(self, file_path: str) -> FileProcessingResult:
        """Generic file processing for unknown types."""
        try:
            path = Path(file_path)
            
            # Try to read as text
            encodings = ['utf-8', 'latin1', 'cp1252']
            content = None
            used_encoding = None
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        content = f.read()
                    used_encoding = encoding
                    break
                except UnicodeDecodeError:
                    continue
            
            if content is None:
                return FileProcessingResult(
                    documents=[], multimodal_elements=[], file_type="unknown",
                    processing_method="binary_file", success=False,
                    error="Could not read file as text"
                )
            
            doc = Document(
                page_content=content,
                metadata={
                    "source": file_path,
                    "filename": path.name,
                    "file_type": path.suffix.lower(),
                    "processing_method": "generic_text",
                    "encoding": used_encoding
                }
            )
            
            return FileProcessingResult(
                documents=[doc],
                multimodal_elements=[],
                file_type="text",
                processing_method="generic_text",
                success=True
            )
            
        except Exception as e:
            return FileProcessingResult(
                documents=[], multimodal_elements=[], file_type="unknown",
                processing_method="generic_text", success=False, error=str(e)
            )
    
    # Helper methods
    
    def _clean_dataframe(self, df: pd.DataFrame) -> pd.DataFrame:
        """Clean and prepare DataFrame for processing."""
        try:
            # Remove empty rows and columns
            df = df.dropna(how='all').dropna(axis=1, how='all')
            
            # Clean column names
            df.columns = [str(col).strip() if col is not None else f"Column_{i}" 
                         for i, col in enumerate(df.columns)]
            
            # Clean cell values
            for col in df.columns:
                if df[col].dtype == 'object':
                    df[col] = df[col].astype(str).str.strip()
                    df[col] = df[col].replace(['nan', 'None', 'null', ''], pd.NA)
            
            return df
            
        except Exception as e:
            logging.warning(f"DataFrame cleaning failed: {e}")
            return df
    
    def _dataframe_to_text(self, df: pd.DataFrame, title: str = "Data Table") -> str:
        """Convert DataFrame to human-readable text."""
        try:
            text_parts = [f"{title}\n" + "=" * len(title)]
            
            # Basic info
            text_parts.append(f"Shape: {df.shape[0]} rows × {df.shape[1]} columns")
            text_parts.append(f"Columns: {', '.join(df.columns[:5])}{'...' if len(df.columns) > 5 else ''}")
            
            # Sample data
            text_parts.append("\nSample data:")
            sample_rows = min(5, len(df))
            for i in range(sample_rows):
                row_data = []
                for col in df.columns[:5]:  # Limit columns
                    value = str(df.iloc[i][col])[:50]  # Limit length
                    row_data.append(f"{col}: {value}")
                text_parts.append(f"Row {i+1}: {', '.join(row_data)}")
            
            if len(df) > sample_rows:
                text_parts.append(f"... and {len(df) - sample_rows} more rows")
            
            return "\n".join(text_parts)
            
        except Exception as e:
            logging.warning(f"DataFrame to text conversion failed: {e}")
            return f"{title}: {df.shape[0]} rows × {df.shape[1]} columns"
    
    def _json_to_text(self, data: Any, indent: int = 0) -> str:
        """Convert JSON data to readable text."""
        try:
            if isinstance(data, dict):
                lines = []
                for key, value in data.items():
                    if isinstance(value, (dict, list)):
                        lines.append(f"{'  ' * indent}{key}:")
                        lines.append(self._json_to_text(value, indent + 1))
                    else:
                        lines.append(f"{'  ' * indent}{key}: {value}")
                return '\n'.join(lines)
            elif isinstance(data, list):
                lines = []
                for i, item in enumerate(data):
                    if isinstance(item, (dict, list)):
                        lines.append(f"{'  ' * indent}Item {i+1}:")
                        lines.append(self._json_to_text(item, indent + 1))
                    else:
                        lines.append(f"{'  ' * indent}- {item}")
                return '\n'.join(lines)
            else:
                return str(data)
        except:
            return json.dumps(data, indent=2)
    
    def _xml_to_text(self, element: ET.Element, indent: int = 0) -> str:
        """Convert XML element to readable text."""
        try:
            lines = []
            
            # Element tag and attributes
            tag_info = element.tag
            if element.attrib:
                attrs = ', '.join(f"{k}={v}" for k, v in element.attrib.items())
                tag_info += f" ({attrs})"
            
            lines.append(f"{'  ' * indent}{tag_info}")
            
            # Element text
            if element.text and element.text.strip():
                lines.append(f"{'  ' * (indent + 1)}{element.text.strip()}")
            
            # Child elements
            for child in element:
                lines.append(self._xml_to_text(child, indent + 1))
            
            return '\n'.join(lines)
        except:
            return ET.tostring(element, encoding='unicode')
    
    def _yaml_to_text(self, data: Any, indent: int = 0) -> str:
        """Convert YAML data to readable text."""
        try:
            return self._json_to_text(data, indent)  # Same structure as JSON
        except:
            return yaml.dump(data, default_flow_style=False)
    
    def get_supported_formats(self) -> Dict[str, Dict[str, Any]]:
        """Get information about all supported file formats."""
        return {
            "documents": {
                "extensions": [".pdf", ".txt", ".docx", ".md", ".rtf"],
                "description": "Text documents and reports",
                "features": ["Text extraction", "Layout analysis", "Metadata"]
            },
            "spreadsheets": {
                "extensions": [".xlsx", ".xls", ".csv"],
                "description": "Spreadsheets and tabular data",
                "features": ["Table extraction", "Multi-sheet support", "Data analysis"],
                "available": EXCEL_AVAILABLE
            },
            "presentations": {
                "extensions": [".pptx"],
                "description": "PowerPoint presentations",
                "features": ["Slide text", "Table extraction", "Image detection"],
                "available": POWERPOINT_AVAILABLE
            },
            "images": {
                "extensions": [".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff", ".webp", ".svg"],
                "description": "Image files",
                "features": ["OCR text extraction", "AI image analysis", "Object detection"],
                "available": PIL_AVAILABLE
            },
            "structured_data": {
                "extensions": [".json", ".xml", ".yaml", ".yml"],
                "description": "Structured data formats",
                "features": ["Data parsing", "Structure preservation", "Table conversion"]
            },
            "web_formats": {
                "extensions": [".html", ".htm"],
                "description": "Web pages and HTML documents",
                "features": ["HTML to text", "Table extraction", "Link preservation"]
            },
            "ebooks": {
                "extensions": [".epub"],
                "description": "Electronic books",
                "features": ["Chapter extraction", "HTML content processing"],
                "available": EBOOK_AVAILABLE
            }
        }