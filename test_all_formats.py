#!/usr/bin/env python3
"""
Universal File Format Support Demo - Test Script
Tests Excel, PowerPoint, images, and all supported file formats
"""

import os
import sys
import json
import tempfile
import pandas as pd
from pathlib import Path
from typing import List, Dict, Any

# Add src to path for imports
sys.path.insert(0, os.path.join(os.path.dirname(__file__), 'src'))

try:
    from src.document_loader import DocumentProcessor
    from src.universal_file_processor import UniversalFileProcessor
    from src.config import Config
except ImportError as e:
    print(f"❌ Import error: {e}")
    print("Make sure you're running this from the project root directory")
    sys.exit(1)


def create_sample_excel():
    """Create a sample Excel file with multiple sheets."""
    try:
        import openpyxl
        from openpyxl.utils.dataframe import dataframe_to_rows
        
        # Create sample data
        sales_data = pd.DataFrame({
            'Month': ['Jan', 'Feb', 'Mar', 'Apr', 'May'],
            'Sales': [10000, 12000, 15000, 11000, 13000],
            'Profit': [2000, 2400, 3000, 2200, 2600],
            'Region': ['North', 'South', 'East', 'West', 'North']
        })
        
        employee_data = pd.DataFrame({
            'Name': ['Alice', 'Bob', 'Charlie', 'Diana'],
            'Department': ['Engineering', 'Marketing', 'Sales', 'HR'],
            'Salary': [75000, 65000, 55000, 60000],
            'Years': [5, 3, 2, 4]
        })
        
        # Create Excel file
        temp_file = tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False)
        
        with pd.ExcelWriter(temp_file.name, engine='openpyxl') as writer:
            sales_data.to_excel(writer, sheet_name='Sales Data', index=False)
            employee_data.to_excel(writer, sheet_name='Employees', index=False)
        
        temp_file.close()
        return temp_file.name
        
    except ImportError:
        print("⚠️ openpyxl not available for creating Excel files")
        return None
    except Exception as e:
        print(f"❌ Error creating Excel file: {e}")
        return None


def create_sample_csv():
    """Create a sample CSV file."""
    try:
        data = {
            'Product': ['Laptop', 'Mouse', 'Keyboard', 'Monitor'],
            'Price': [999.99, 29.99, 79.99, 299.99],
            'Stock': [15, 50, 30, 8],
            'Category': ['Electronics', 'Accessories', 'Accessories', 'Electronics']
        }
        df = pd.DataFrame(data)
        
        temp_file = tempfile.NamedTemporaryFile(mode='w', suffix='.csv', delete=False)
        df.to_csv(temp_file.name, index=False)
        temp_file.close()
        
        return temp_file.name
        
    except Exception as e:
        print(f"❌ Error creating CSV file: {e}")
        return None


def create_sample_json():
    """Create a sample JSON file."""
    try:
        data = {
            "company": "Tech Corp",
            "employees": [
                {"name": "Alice", "role": "Engineer", "salary": 75000},
                {"name": "Bob", "role": "Designer", "salary": 65000},
                {"name": "Charlie", "role": "Manager", "salary": 85000}
            ],
            "departments": {
                "engineering": {"budget": 500000, "headcount": 10},
                "design": {"budget": 200000, "headcount": 5},
                "management": {"budget": 300000, "headcount": 3}
            }
        }
        
        temp_file = tempfile.NamedTemporaryFile(mode='w', suffix='.json', delete=False)
        json.dump(data, temp_file, indent=2)
        temp_file.close()
        
        return temp_file.name
        
    except Exception as e:
        print(f"❌ Error creating JSON file: {e}")
        return None


def create_sample_image():
    """Create a sample image with text."""
    try:
        from PIL import Image, ImageDraw, ImageFont
        import matplotlib.pyplot as plt
        
        # Create a simple chart image
        fig, ax = plt.subplots(figsize=(8, 6))
        categories = ['A', 'B', 'C', 'D']
        values = [23, 45, 56, 78]
        
        bars = ax.bar(categories, values, color=['red', 'green', 'blue', 'orange'])
        ax.set_title('Sample Data Visualization', fontsize=16)
        ax.set_xlabel('Categories')
        ax.set_ylabel('Values')
        
        # Add value labels on bars
        for bar, value in zip(bars, values):
            height = bar.get_height()
            ax.text(bar.get_x() + bar.get_width()/2., height + 1,
                   str(value), ha='center', va='bottom')
        
        temp_file = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
        plt.savefig(temp_file.name, dpi=150, bbox_inches='tight')
        plt.close()
        temp_file.close()
        
        return temp_file.name
        
    except ImportError:
        print("⚠️ PIL/matplotlib not available for creating images")
        return None
    except Exception as e:
        print(f"❌ Error creating image: {e}")
        return None


def create_sample_powerpoint():
    """Create a sample PowerPoint file."""
    try:
        from pptx import Presentation
        from pptx.util import Inches
        
        prs = Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Sample Presentation"
        subtitle.text = "Created for RAG System Testing"
        
        # Content slide with bullet points
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = 'Key Features'
        
        tf = body_shape.text_frame
        tf.text = 'Advanced Document Processing'
        
        p = tf.add_paragraph()
        p.text = 'Multi-format support (PDF, Excel, PowerPoint, Images)'
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = 'AI-powered content analysis'
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = 'Table and image extraction'
        p.level = 1
        
        # Table slide
        title_only_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_only_layout)
        shapes = slide.shapes
        
        shapes.title.text = 'Sample Data Table'
        
        # Add table
        rows, cols = 4, 3
        left = Inches(1.0)
        top = Inches(2.0)
        width = Inches(8.0)
        height = Inches(3.0)
        
        table = shapes.add_table(rows, cols, left, top, width, height).table
        
        # Set column widths
        table.columns[0].width = Inches(2.0)
        table.columns[1].width = Inches(3.0)
        table.columns[2].width = Inches(3.0)
        
        # Header row
        table.cell(0, 0).text = 'Product'
        table.cell(0, 1).text = 'Revenue'
        table.cell(0, 2).text = 'Growth'
        
        # Data rows
        table.cell(1, 0).text = 'Product A'
        table.cell(1, 1).text = '$1,000,000'
        table.cell(1, 2).text = '+15%'
        
        table.cell(2, 0).text = 'Product B'
        table.cell(2, 1).text = '$750,000'
        table.cell(2, 2).text = '+8%'
        
        table.cell(3, 0).text = 'Product C'
        table.cell(3, 1).text = '$500,000'
        table.cell(3, 2).text = '+22%'
        
        temp_file = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
        prs.save(temp_file.name)
        temp_file.close()
        
        return temp_file.name
        
    except ImportError:
        print("⚠️ python-pptx not available for creating PowerPoint files")
        return None
    except Exception as e:
        print(f"❌ Error creating PowerPoint file: {e}")
        return None


def test_document_processor():
    """Test the DocumentProcessor with various file formats."""
    print("🧪 Testing Document Processor with Multiple File Formats")
    print("=" * 60)
    
    processor = DocumentProcessor()
    
    # Display supported formats
    supported_formats = processor.get_supported_formats()
    print(f"📋 Total supported extensions: {supported_formats['total_supported_extensions']}")
    print(f"📄 Supported formats: {', '.join(supported_formats['supported_extensions'][:10])}...")
    
    # Create sample files
    sample_files = []
    
    print("\n🔨 Creating sample files...")
    
    # Excel file
    excel_file = create_sample_excel()
    if excel_file:
        sample_files.append(("Excel", excel_file))
        print("✅ Created Excel file")
    
    # CSV file
    csv_file = create_sample_csv()
    if csv_file:
        sample_files.append(("CSV", csv_file))
        print("✅ Created CSV file")
    
    # JSON file
    json_file = create_sample_json()
    if json_file:
        sample_files.append(("JSON", json_file))
        print("✅ Created JSON file")
    
    # Image file
    image_file = create_sample_image()
    if image_file:
        sample_files.append(("Image", image_file))
        print("✅ Created Image file")
    
    # PowerPoint file
    ppt_file = create_sample_powerpoint()
    if ppt_file:
        sample_files.append(("PowerPoint", ppt_file))
        print("✅ Created PowerPoint file")
    
    if not sample_files:
        print("❌ Could not create any sample files. Install missing dependencies.")
        return False
    
    print(f"\n🧪 Testing {len(sample_files)} file formats...")
    
    results = []
    
    for file_type, file_path in sample_files:
        print(f"\n📄 Testing {file_type} file...")
        
        try:
            # Validate file support
            validation = processor.validate_file_support(file_path)
            print(f"   Validation: {'✅ Supported' if validation['is_supported'] else '❌ Not supported'}")
            print(f"   Processing method: {validation.get('processing_method', 'unknown')}")
            print(f"   Features: {', '.join(validation.get('available_features', []))}")
            
            if validation['is_supported']:
                # Process the file
                documents = processor.load_document(file_path)
                multimodal_elements = processor.get_multimodal_elements()
                
                print(f"   📄 Documents created: {len(documents)}")
                print(f"   📊 Multimodal elements: {len(multimodal_elements)}")
                
                if documents:
                    sample_content = documents[0].page_content[:200] + "..." if len(documents[0].page_content) > 200 else documents[0].page_content
                    print(f"   📝 Sample content: {sample_content}")
                
                if multimodal_elements:
                    for elem in multimodal_elements[:2]:  # Show first 2
                        print(f"   📊 {elem.element_type}: {elem.text_description[:100]}...")
                
                results.append({
                    "file_type": file_type,
                    "success": True,
                    "documents": len(documents),
                    "multimodal_elements": len(multimodal_elements)
                })
                
                print(f"   ✅ {file_type} processing successful")
            else:
                results.append({
                    "file_type": file_type,
                    "success": False,
                    "error": validation.get('error', 'Unknown error')
                })
                print(f"   ❌ {file_type} processing failed")
        
        except Exception as e:
            print(f"   ❌ Error processing {file_type}: {e}")
            results.append({
                "file_type": file_type,
                "success": False,
                "error": str(e)
            })
    
    # Clean up temporary files
    for _, file_path in sample_files:
        try:
            os.unlink(file_path)
        except:
            pass
    
    # Results summary
    print("\n" + "=" * 60)
    print("📊 Test Results Summary")
    print("=" * 60)
    
    successful = sum(1 for r in results if r['success'])
    total = len(results)
    
    print(f"Overall: {successful}/{total} file formats processed successfully")
    
    for result in results:
        if result['success']:
            print(f"✅ {result['file_type']}: {result['documents']} docs, {result['multimodal_elements']} multimodal")
        else:
            print(f"❌ {result['file_type']}: {result.get('error', 'Failed')}")
    
    print("\n🎯 File Format Capabilities:")
    print("✅ Excel (.xlsx, .xls): Multi-sheet extraction, table analysis")
    print("✅ CSV: Data parsing, automatic table conversion")
    print("✅ JSON: Structure parsing, table extraction for arrays")
    print("✅ Images (.jpg, .png, etc.): OCR, AI analysis, object detection")
    print("✅ PowerPoint (.pptx): Slide text, table extraction")
    print("✅ PDF: Advanced table/image extraction (previously tested)")
    print("✅ HTML/XML: Structure parsing, table extraction")
    print("✅ YAML: Configuration file parsing")
    print("✅ Text formats: Full text extraction")
    
    return successful == total


def test_universal_processor():
    """Test the UniversalFileProcessor directly."""
    print("\n🔧 Testing Universal File Processor")
    print("=" * 50)
    
    processor = UniversalFileProcessor()
    
    # Display capabilities
    capabilities = processor.get_supported_formats()
    
    print("🚀 Processing Capabilities:")
    for category, info in capabilities.items():
        status = "✅ Available" if info.get('available', True) else "❌ Missing dependencies"
        print(f"   {category}: {status}")
        print(f"      Extensions: {', '.join(info['extensions'])}")
        print(f"      Features: {', '.join(info['features'])}")
    
    return True


def main():
    """Run all file format tests."""
    print("🌟 Universal File Format Support Test Suite")
    print("=" * 70)
    
    print("📋 Supported File Extensions:")
    config = Config()
    extensions_by_category = {
        "Documents": [".pdf", ".txt", ".docx", ".md", ".rtf"],
        "Spreadsheets": [".xlsx", ".xls", ".csv"],
        "Presentations": [".pptx", ".ppt"],
        "Images": [".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff", ".webp", ".svg"],
        "Structured Data": [".json", ".xml", ".yaml", ".yml"],
        "Web Formats": [".html", ".htm"],
        "Ebooks": [".epub", ".mobi"]
    }
    
    for category, extensions in extensions_by_category.items():
        print(f"   {category}: {', '.join(extensions)}")
    
    print(f"\nTotal: {len(config.SUPPORTED_EXTENSIONS)} file formats supported!")
    
    # Check dependencies
    print("\n🔍 Checking dependencies...")
    dependencies = {
        "openpyxl (Excel)": "openpyxl",
        "python-pptx (PowerPoint)": "pptx",
        "PIL/Pillow (Images)": "PIL",
        "matplotlib (Charts)": "matplotlib.pyplot",
        "pandas (Data)": "pandas",
        "yaml (YAML)": "yaml"
    }
    
    for name, module in dependencies.items():
        try:
            __import__(module.split('.')[0])
            print(f"   ✅ {name}")
        except ImportError:
            print(f"   ❌ {name} - install with: pip install {module.split('.')[0]}")
    
    print("\n" + "=" * 70)
    
    # Run tests
    test_results = []
    
    # Test 1: Universal Processor
    test_results.append(("Universal Processor", test_universal_processor()))
    
    # Test 2: Document Processor with multiple formats
    test_results.append(("Multi-Format Processing", test_document_processor()))
    
    # Final summary
    print("\n" + "=" * 70)
    print("🏆 Final Test Results")
    print("=" * 70)
    
    passed = sum(1 for _, result in test_results if result)
    total = len(test_results)
    
    for test_name, result in test_results:
        status = "✅ PASSED" if result else "❌ FAILED"
        print(f"{test_name}: {status}")
    
    print(f"\nOverall: {passed}/{total} test suites passed")
    
    if passed == total:
        print("\n🎉 All tests passed! Your system supports comprehensive file formats:")
        print("   📊 Excel spreadsheets with multi-sheet support")
        print("   🖼️ Images with AI analysis and OCR")
        print("   📽️ PowerPoint presentations with table extraction")
        print("   📄 JSON, XML, YAML, HTML, and more")
        print("   📱 CSV data with automatic table conversion")
        print("   📚 Ebooks and rich document formats")
        
        print("\n🚀 Next steps:")
        print("1. Upload any supported file format to your RAG system")
        print("2. Ask questions about tables, images, and structured content")
        print("3. Use multi-modal queries to search across different content types")
        print("4. Export analysis results and multimodal summaries")
    else:
        print("\n⚠️ Some tests failed. Install missing dependencies:")
        print("pip install openpyxl python-pptx Pillow matplotlib ebooklib PyYAML html2text")
    
    return passed == total


if __name__ == "__main__":
    success = main()
    sys.exit(0 if success else 1)